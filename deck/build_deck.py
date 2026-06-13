#!/usr/bin/env python3
"""Fill the official KSP Datathon submission template with TRINETRA content.

Keeps the KSP branding (header/footer) and per-slide titles intact; adds body
content, two shape diagrams (process flow + architecture), and the four live
prototype screenshots. Output: TRINETRA_Submission.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = "/Users/raymac/Downloads/KSP Datathon 2026 _ Prototype Submission Template.pptx"
SHOTS = os.path.join(HERE, "shots")
OUT = os.path.join(HERE, "TRINETRA_Submission.pptx")

# palette
INK = RGBColor(0x1B, 0x26, 0x34)
MUTED = RGBColor(0x5A, 0x6B, 0x7B)
RED = RGBColor(0xC0, 0x39, 0x2B)
BLUE = RGBColor(0x1C, 0x6F, 0xB0)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
AMBER = RGBColor(0xB9, 0x7A, 0x12)
LIGHT = RGBColor(0xED, 0xF2, 0xF6)
CARD = RGBColor(0xF4, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD4, 0xDD, 0xE6)

prs = Presentation(TEMPLATE)
slides = list(prs.slides)


def title_of(slide):
    cands = [s for s in slide.shapes if s.has_text_frame and s.text.strip()]
    return min(cands, key=lambda s: s.top) if cands else None


def content_top(slide, gap=0.22):
    t = title_of(slide)
    if t is None:
        return 1.5
    return (t.top + t.height) / 914400 + gap


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    return tb, tf


def para(tf, text, size=10.5, bold=False, color=INK, first=False,
         bullet=False, space_after=4, level=0, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    runs = text if isinstance(text, list) else [(text, color, bold)]
    for i, item in enumerate(runs):
        txt, col, bd = item
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = bd
        r.font.color.rgb = col
        r.font.name = "Calibri"
    if bullet:
        _bullet(p)
    return p


def _bullet(p):
    pPr = p._p.get_or_add_pPr()
    buFont = pPr.makeelement(qn('a:buFont'), {'typeface': 'Arial'})
    buChar = pPr.makeelement(qn('a:buChar'), {'char': '▪'})
    pPr.append(buFont); pPr.append(buChar)
    pPr.set('indent', str(-Inches(0.18).emu))
    pPr.set('marL', str(Inches(0.18).emu))


def rect(slide, x, y, w, h, fill=CARD, line=LINE, rounded=True, line_w=1.0):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = 0.08
    except Exception:
        pass
    return shp


def shape_text(shp, lines, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, (txt, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(1); p.space_before = Pt(0)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = "Calibri"


def arrow(slide, x, y, w=0.32, h=0.22, color=RED):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background(); a.shadow.inherit = False
    return a


def down_arrow(slide, x, y, w=0.22, h=0.3, color=RED):
    a = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background(); a.shadow.inherit = False
    return a


def pic(slide, path, x, y, w):
    return slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))


# ----------------------------------------------------------------- Slide 1
s = slides[0]
box = title_of(s)
tf = box.text_frame
tf.clear()
para(tf, [("Team Details", INK, True)], size=15, first=True, space_after=8)
fields = [
    ("Team name: ", "Trinetra Labs"),
    ("Team leader name: ", "Rayhankhan Pathan"),
    ("Team size: ", "1"),
    ("Problem Statement: ", "PS1 — Intelligent Conversational AI for the KSP Crime Database"),
]
for k, v in fields:
    para(tf, [(k, INK, True), (v, MUTED, False)], size=11, space_after=5)

# ----------------------------------------------------------------- Slide 2
s = slides[1]
ct = content_top(s)
_, tf = textbox(s, 0.4, ct, 9.2, 1.2)
para(tf, [("TRINETRA", RED, True), (" — an agentic crime-intelligence console for the Karnataka State Police. An investigator asks a question in Kannada or English (voice or text); the AI agent queries the crime database, uncovers hidden criminal networks, profiles repeat offenders, and answers with an interactive graph, hotspot map and timeline — every fact cited to its FIR.", INK, False)], size=12, first=True, space_after=6)
# three stat cards
cards = [("Days → seconds", "FIR-to-lead cross-referencing"),
         ("EN + ಕನ್ನಡ", "voice & text, context-aware"),
         ("100% on Catalyst", "deployed, explainable, secure")]
cw, gap = 2.93, 0.2
x0 = 0.4
cy = ct + 1.5
for i, (big, lab) in enumerate(cards):
    x = x0 + i * (cw + gap)
    r = rect(s, x, cy, cw, 1.15, fill=CARD, line=LINE)
    shape_text(r, [(big, 15, True, RED), (lab, 9.5, False, MUTED)])

# ----------------------------------------------------------------- Slide 3
s = slides[2]
ct = content_top(s, 0.18)
cols = [
    ("How it's different", BLUE,
     ["Not a dashboard and not a plain chatbot — an agent that returns understanding.",
      "Generative UI: it decides whether to draw a network, a map or a timeline.",
      "Relationships modelled as edge tables → graph in code (no graph DB needed).",
      "Kannada + English voice, built for Indian policing."]),
    ("How it solves the problem", TEAL,
     ["Natural language → safe ZCQL over the Data Store, no SQL skills needed.",
      "Surfaces hidden offender / phone / money links across jurisdictions.",
      "Transparent risk scoring prioritises investigation.",
      "Cuts manual cross-referencing from days to seconds."]),
    ("USP", RED,
     ["Agentic + explainable: every answer cites its FIRs (evidence trail).",
      "Spatiotemporal hotspots + organised-network detection in one console.",
      "Fully on Zoho Catalyst — Qwen LLM, Zia voice, Data Store, Slate.",
      "Production-grade: role-based access, audit-ready, scalable."]),
]
cw, gap = 3.0, 0.1
x0 = 0.35
cy = max(ct, 2.05)
for i, (head, col, items) in enumerate(cols):
    x = x0 + i * (cw + gap)
    r = rect(s, x, cy, cw, 3.0, fill=WHITE, line=LINE)
    tf = r.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(7); tf.margin_right = Pt(7); tf.margin_top = Pt(7)
    para(tf, [(head, col, True)], size=11.5, first=True, space_after=6)
    for it in items:
        para(tf, it, size=9.3, color=INK, bullet=True, space_after=5)

# ----------------------------------------------------------------- Slide 4
s = slides[3]
ct = content_top(s)
feats = [
    ("Conversational agent", "Voice + text, English & Kannada, context-aware follow-ups."),
    ("NL → ZCQL retrieval", "Plain questions become safe, validated Data Store queries."),
    ("Criminal network analysis", "Hidden links between offenders, victims, phones, accounts."),
    ("Offender risk profiling", "Explainable 0–100 score from weighted factors."),
    ("Spatiotemporal hotspots", "District clusters layered with time-of-day."),
    ("Investigator timeline", "Auto case chronology showing MO escalation."),
    ("Money-trail analysis", "Suspicious transactions and mule-account links."),
    ("Similar-case leads", "Find look-alike cases by MO — investigative leads."),
    ("Crime forecasting", "Projected next-week hotspots for proactive patrol."),
    ("Early-warning alerts", "Spikes, forming rings, repeat offenders."),
    ("Explainable & cited", "Shows the generated ZCQL + reasoning; cites FIRs; PDF export."),
    ("Secure & governed", "Role-based access, audit logs, SELECT-only query gate."),
]
cw, ch, gx, gy = 4.55, 0.54, 0.1, 0.06
x0 = 0.4
cy = ct
for i, (h, d) in enumerate(feats):
    col = i % 2; row = i // 2
    x = x0 + col * (cw + gx); y = cy + row * (ch + gy)
    r = rect(s, x, y, cw, ch, fill=CARD, line=LINE)
    tf = r.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(8); tf.margin_right = Pt(6); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    para(tf, [(h + "  ", INK, True)], size=10, first=True, space_after=1)
    para(tf, [(d, MUTED, False)], size=8.6, space_after=0)

# ----------------------------------------------------------------- Slide 5 (process flow)
s = slides[4]
ct = content_top(s)
flow = [
    ("Investigator", "asks in EN / ಕನ್ನಡ\n(voice or text)", BLUE),
    ("Agent", "intent → generate ZCQL\n→ safety gate", RED),
    ("Data Store", "validated ZCQL\nover crime DB", TEAL),
    ("Analysis", "network · risk ·\nhotspot · money", AMBER),
    ("Answer", "graph / map / timeline\n+ cited FIRs · PDF", BLUE),
]
n = len(flow); bw = 1.62; bh = 1.0; gap = 0.21
total = n * bw + (n - 1) * gap
x0 = (10 - total) / 2
fy = max(ct + 0.35, 2.1)
for i, (h, d, col) in enumerate(flow):
    x = x0 + i * (bw + gap)
    r = rect(s, x, fy, bw, bh, fill=WHITE, line=col, line_w=1.5)
    shape_text(r, [(h, 11, True, col)] + [(ln, 8.3, False, MUTED) for ln in d.split("\n")])
    if i < n - 1:
        arrow(s, x + bw + (gap - 0.32) / 2, fy + bh / 2 - 0.11)
_, tf = textbox(s, 0.4, fy + bh + 0.35, 9.2, 0.7)
para(tf, [("Use case: ", RED, True), ("“Who connects these chain-snatching cases?” → the agent runs the query, finds 4 FIRs tied to one offender via a shared phone + mule account, and renders the ring with an evidence trail — in seconds.", INK, False)], size=10, first=True)

# ----------------------------------------------------------------- Slide 6 (wireframe / mock)
s = slides[5]
ct = content_top(s)
img_w = 5.7
p = pic(s, os.path.join(SHOTS, "network.png"), 0.4, ct, img_w)
r = rect(s, 6.35, ct, 3.25, 3.3, fill=CARD, line=LINE)
tf = r.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
tf.margin_left = Pt(9); tf.margin_right = Pt(8); tf.margin_top = Pt(8)
para(tf, [("Console layout", INK, True)], size=11.5, first=True, space_after=6)
for h, d in [("Command bar", "natural-language + voice query"),
             ("Network canvas", "force-directed graph, neighbour highlight"),
             ("Inspector", "risk factors + connected entities"),
             ("Agent log", "cited answer · view query & reasoning · PDF")]:
    para(tf, [(h + " — ", RED, True), (d, MUTED, False)], size=9.2, bullet=True, space_after=6)

# ----------------------------------------------------------------- Slide 7 (architecture)
s = slides[6]
ct = content_top(s)
ay = max(ct + 0.1, 1.85)
# layer 1: frontend
l1 = rect(s, 0.7, ay, 8.6, 0.62, fill=RGBColor(0x16,0x24,0x3A), line=None)
shape_text(l1, [[("Frontend — Catalyst Slate (Next.js console: chat · network graph · hotspot map · timeline)", 10.5, True, WHITE)]][0])
down_arrow(s, 4.9, ay + 0.64, color=MUTED)
# layer 2: gateway + functions
l2 = rect(s, 0.7, ay + 1.0, 8.6, 0.62, fill=RGBColor(0x1C,0x6F,0xB0), line=None)
shape_text(l2, [("API Gateway  →  Functions: agent (NL→ZCQL) · network-build", 10.5, True, WHITE)])
down_arrow(s, 4.9, ay + 1.64, color=MUTED)
# layer 3: services row
svcs = [("QuickML\nQwen LLM + RAG", TEAL), ("Zia\nvoice EN/KN", AMBER),
        ("Data Store\nZCQL + edges", BLUE), ("Stratus\nfiles", RED),
        ("SmartBrowz\nPDF", TEAL), ("Auth\nRBAC + audit", RGBColor(0x6D,0x4A,0xA0))]
sw = 1.42; sg = 0.05; sx = (10 - (len(svcs)*sw + (len(svcs)-1)*sg)) / 2
sy = ay + 2.05
for i, (t, col) in enumerate(svcs):
    x = sx + i*(sw+sg)
    r = rect(s, x, sy, sw, 0.85, fill=WHITE, line=col, line_w=1.5)
    lines = t.split("\n")
    shape_text(r, [(lines[0], 9.3, True, col), (lines[1], 8, False, MUTED)])
_, tf = textbox(s, 0.7, sy + 1.0, 8.6, 0.5)
para(tf, [("Synthetic data → Data Store. Every generated ZCQL is validated (SELECT-only, whitelisted tables) before execution — the governance gate.", MUTED, False)], size=8.8, first=True, align=PP_ALIGN.CENTER)

# ----------------------------------------------------------------- Slide 8 (technologies)
s = slides[7]
ct = content_top(s)
groups = [
    ("Frontend", BLUE, "Next.js 14 · React · TypeScript · Tailwind CSS · react-force-graph-2d · MapLibre GL · Web Speech API"),
    ("Backend / Functions", TEAL, "Node.js Catalyst Functions (Advanced I/O) · ZCQL · REST + API Gateway"),
    ("AI / ML", RED, "Qwen via Catalyst QuickML LLM-Serving · RAG knowledge base · Zia STT / TTS / translation"),
    ("Data", AMBER, "Catalyst Data Store (relational + edge tables) · Python synthetic-data generator (planted rings)"),
    ("Platform", RGBColor(0x6D,0x4A,0xA0), "Zoho Catalyst — Slate hosting · Authentication · Stratus · SmartBrowz · Cron · Signals"),
]
cy = ct
for i, (h, col, d) in enumerate(groups):
    y = cy + i * 0.72
    bar = rect(s, 0.4, y, 0.12, 0.6, fill=col, line=None, rounded=False)
    _, tf = textbox(s, 0.65, y, 8.9, 0.66, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(h, col, True)], size=11, first=True, space_after=1)
    para(tf, [(d, INK, False)], size=9.5)

# ----------------------------------------------------------------- Slide 9 (Catalyst services)
s = slides[8]
ct = content_top(s)
svc_rows = [
    ("Slate / Web Hosting", "Hosts the Next.js console"),
    ("Functions (Advanced I/O)", "Agent (NL→ZCQL) & network-build logic"),
    ("API Gateway", "Routing, throttling, auth in front of Functions"),
    ("QuickML — LLM Serving + RAG", "Qwen for NL→ZCQL & grounded answers"),
    ("Zia Services", "Kannada/English speech-to-text, TTS, translation"),
    ("Data Store (ZCQL)", "Crime DB + relationship edge tables + search"),
    ("Stratus", "Case files, evidence, exported PDFs"),
    ("SmartBrowz", "Conversation → PDF report generation"),
    ("Authentication", "Role-based access (investigator/analyst/supervisor)"),
    ("Cron + Signals", "Scheduled risk recompute & early-warning alerts"),
]
cw, ch, gx, gy = 4.55, 0.6, 0.1, 0.06
x0 = 0.4
cy = ct
for i, (h, d) in enumerate(svc_rows):
    col = i % 2; row = i // 2
    x = x0 + col * (cw + gx); y = cy + row * (ch + gy)
    r = rect(s, x, y, cw, ch, fill=CARD, line=LINE)
    tf = r.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(8); tf.margin_right = Pt(6)
    para(tf, [(h, TEAL, True)], size=9.7, first=True, space_after=1)
    para(tf, [(d, MUTED, False)], size=8.4)

# ----------------------------------------------------------------- Slide 10 (cost)
s = slides[9]
ct = content_top(s)
_, tf = textbox(s, 0.4, ct, 9.2, 0.5)
para(tf, [("Built within the hackathon Catalyst credits (₹1,500 + $250). Indicative monthly run-rate at SCRB pilot scale:", INK, False)], size=10.5, first=True)
rows = [("Catalyst service", "Usage", "Indicative / month"),
        ("Functions (agent + network)", "~50k invocations", "Within free tier → low"),
        ("Data Store", "~2–5 lakh records", "Low"),
        ("QuickML LLM-Serving (Qwen)", "~30k queries", "Usage-based"),
        ("Zia voice", "~5k minutes", "Usage-based"),
        ("Slate + Stratus + SmartBrowz", "hosting + PDFs", "Low")]
ty = ct + 0.65; rh = 0.46
c1, c2, c3 = 4.0, 2.6, 2.6
tx = 0.4
for i, (a, b, c) in enumerate(rows):
    head = i == 0
    fill = RGBColor(0x16,0x24,0x3A) if head else (CARD if i % 2 else WHITE)
    for cx, cwid, val, al in [(tx, c1, a, PP_ALIGN.LEFT), (tx+c1, c2, b, PP_ALIGN.LEFT), (tx+c1+c2, c3, c, PP_ALIGN.LEFT)]:
        cell = rect(s, cx, ty + i*rh, cwid, rh, fill=fill, line=LINE, rounded=False)
        shape_text(cell, [(val, 9.3, head, WHITE if head else INK)], align=al, anchor=MSO_ANCHOR.MIDDLE)
        cell.text_frame.margin_left = Pt(8)

# ----------------------------------------------------------------- Slide 11 (snapshots)
s = slides[10]
ct = content_top(s, 0.12)
imgW = 2.75; imgH = imgW / (1366/820)
shots = [("network.png", "Network — organised ring"),
         ("map.png", "City hotspot map (MapLibre)"),
         ("profiles.png", "Offender risk profiles"),
         ("timeline.png", "Case timeline")]
colx = [1.95, 5.30]
rowy = [max(ct, 1.5), max(ct, 1.5) + imgH + 0.33]
for i, (fn, cap) in enumerate(shots):
    cidx = i % 2; ridx = i // 2
    x = colx[cidx]; y = rowy[ridx]
    _, tf = textbox(s, x, y - 0.02, imgW, 0.2)
    para(tf, [(cap, MUTED, True)], size=8.6, first=True)
    pic(s, os.path.join(SHOTS, fn), x, y + 0.18, imgW)

# ----------------------------------------------------------------- Slide 12 (performance)
s = slides[11]
ct = content_top(s)
stats = [("< 1s", "agent response (mock mode)"),
         ("100%", "malicious queries blocked by ZCQL gate"),
         ("1.5 lakh", "FIRs generated & queried in tests"),
         ("6/6", "demo intents return grounded, cited answers"),
         ("16 nodes", "network built per offender ego-graph"),
         ("11 tables", "relational schema, edge-modelled")]
cw, chh, gx, gy = 2.93, 1.15, 0.2, 0.18
x0 = 0.4
cy = ct + 0.05
for i, (big, lab) in enumerate(stats):
    col = i % 3; row = i // 3
    x = x0 + col*(cw+gx); y = cy + row*(chh+gy)
    r = rect(s, x, y, cw, chh, fill=CARD, line=LINE)
    shape_text(r, [(big, 26, True, RED), (lab, 9, False, MUTED)])

# ----------------------------------------------------------------- Slide 13 (links)
s = slides[12]
ct = content_top(s, 0.25)
links = [("GitHub repository", "https://github.com/mrayhankhan/trinetra", BLUE),
         ("Demo video (3 min)", "https://youtu.be/<unlisted-link>", RED),
         ("Deployed link (Catalyst)", "https://trinetra-60074156791.development.catalystserverless.in/app/index.html", TEAL)]
cy = max(ct, 2.95)
for i, (h, url, col) in enumerate(links):
    y = cy + i*0.72
    r = rect(s, 0.5, y, 9.0, 0.58, fill=CARD, line=LINE)
    tf = r.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(12)
    para(tf, [(h + ":   ", col, True), (url, INK, False)], size=11.5, first=True)
_, tf = textbox(s, 0.5, cy + 3*0.72 + 0.02, 9.0, 0.35)
para(tf, [("Replace with your public links before submitting. Deployment must be on Catalyst.", MUTED, False)], size=9, first=True)

# ----------------------------------------------------------------- Slide 14 (future)
s = slides[13]
ct = content_top(s)
fut = [
    ("RAG criminology knowledge base", "Ground answers in IPC / case-law / SOP documents via QuickML RAG."),
    ("Crime forecasting", "QuickML time-series + Zia AutoML to predict emerging hotspots."),
    ("Financial-crime expansion", "Deeper money-mule and cross-bank transaction graphs."),
    ("More Indic languages", "Beyond Kannada — Hindi, Tamil, Telugu, Urdu."),
    ("Mobile + field app", "Push alerts and on-the-go queries for officers."),
    ("Live SCRB integration", "Connect to real KSP systems with full governance & DPDP compliance."),
]
cy = ct
for i, (h, d) in enumerate(fut):
    y = cy + i*0.6
    bar = rect(s, 0.4, y + 0.05, 0.12, 0.42, fill=RED, line=None, rounded=False)
    _, tf = textbox(s, 0.65, y, 8.9, 0.58, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(h + " — ", INK, True), (d, MUTED, False)], size=10, first=True)

# ----------------------------------------------------------------- Cleanup
# The template's last slide is already a clean "THANK YOU" close, so we add
# nothing to it. Drop the template's spare "Blank slide" (slide 15).
xml_slides = prs.slides._sldIdLst
xml_slides.remove(list(xml_slides)[14])

prs.save(OUT)
print("saved", OUT)
